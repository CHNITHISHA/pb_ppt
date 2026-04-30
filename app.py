import pandas as pd
from pptx import Presentation

# -------------------------------
# STEP 1: READ CSV FILE
# -------------------------------
file_path = "graph1.csv"   # your exported file

df = pd.read_csv(file_path)

print("\n=== RAW DATA ===")
print(df)


# -------------------------------
# STEP 2: CLEAN COLUMN NAMES
# -------------------------------
# Remove spaces if any (Power BI sometimes adds spaces)
df.columns = df.columns.str.strip()

print("\nColumns:", df.columns)


# -------------------------------
# STEP 3: CHECK REQUIRED COLUMNS
# -------------------------------
if "Product" not in df.columns or "Sales" not in df.columns:
    print("\n❌ ERROR: Column names mismatch!")
    print("Available columns:", df.columns)
    exit()


# -------------------------------
# STEP 4: PROCESS DATA (OPTIONAL)
# -------------------------------
# If duplicate products exist → group them
df = df.groupby("Product", as_index=False)["Sales"].sum()

print("\n=== PROCESSED DATA ===")
print(df)


# -------------------------------
# STEP 5: LOAD PPT TEMPLATE
# -------------------------------
ppt_file = "powerbi_template.pptx"
prs = Presentation(ppt_file)

# Slide 2 (index starts from 0)
slide = prs.slides[1]


# -------------------------------
# STEP 6: FIND TABLE IN SLIDE
# -------------------------------
table = None

for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    raise Exception("❌ No table found in Slide 2")


# -------------------------------
# STEP 7: CLEAR OLD DATA
# -------------------------------
for i in range(1, len(table.rows)):
    for j in range(len(table.columns)):
        table.cell(i, j).text = ""


# -------------------------------
# STEP 8: INSERT NEW DATA
# -------------------------------
for i, row in df.iterrows():
    if i + 1 >= len(table.rows):
        print("⚠️ Not enough rows in PPT table")
        break

    table.cell(i + 1, 0).text = str(row["Product"])
    table.cell(i + 1, 1).text = str(row["Sales"])


# -------------------------------
# STEP 9: SAVE OUTPUT PPT
# -------------------------------
output_file = "final_output.pptx"
prs.save(output_file)

print(f"\n✅ SUCCESS: PPT updated → {output_file}")
